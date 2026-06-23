#!/usr/bin/env python3
"""
PPTX 风格数据提取库 (v2)
基于 XML 直接解析，输出精准 Markdown 结构报告。
不依赖 python-pptx / Pillow，仅使用 zipfile + xml.etree。

用法:
    from scripts.parse_pptx import parse_pptx_to_markdown
    md_text, image_map = parse_pptx_to_markdown("/path/to/file.pptx", "/tmp/output")

    # 带图片资源公开 URL 替换:
    md_text, image_map = parse_pptx_to_markdown(
        "/path/to/file.pptx", "/tmp/output",
        resource_base_url="https://bucket.oss-cn-hangzhou.aliyuncs.com/user/.../style/task123"
    )
"""

import os
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

# ============================================================
# XML Namespaces
# ============================================================
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

COLOR_ROLE_DESCRIPTIONS = {
    "dk1": "主要文字颜色",
    "lt1": "页面背景颜色",
    "dk2": "次级文字 / 深色强调",
    "lt2": "辅助浅色背景",
    "accent1": "主强调色 / 品牌色",
    "accent2": "图表分类色 2",
    "accent3": "图表分类色 3",
    "accent4": "图表分类色 4",
    "accent5": "图表分类色 5",
    "accent6": "图表分类色 6",
}


# ============================================================
# XML Helpers
# ============================================================
def q(tag: str) -> str:
    """Convert 'ns:tag' or XPath 'ns:a/ns:b' to Clark notation."""
    if "}" in tag:
        return tag
    return "/".join(f"{{{NS[ns]}}}{name}" for ns, name in (p.split(":") for p in tag.split("/")))


def load_xml(path: Path) -> ET.Element | None:
    if not path.exists():
        return None
    return ET.parse(str(path)).getroot()


def load_rels(path: Path) -> dict[str, dict[str, str]]:
    rels = {}
    root = load_xml(path)
    if root is None:
        return rels
    for rel in root.findall(q("rel:Relationship")):
        rels[rel.get("Id")] = {"Target": rel.get("Target"), "Type": rel.get("Type")}
    return rels


def emu_to_pt(emu: str | None) -> float | None:
    try:
        return round(int(emu) / 12700, 2)
    except (TypeError, ValueError):
        return None


def resolve_color(color_elem: ET.Element | None) -> str | None:
    if color_elem is None:
        return None
    srgb = color_elem.find(q("a:srgbClr"))
    if srgb is not None:
        val = srgb.get("val")
        if val:
            return f"#{val}"
    scheme = color_elem.find(q("a:schemeClr"))
    if scheme is not None:
        return f"scheme:{scheme.get('val')}"
    sys_clr = color_elem.find(q("a:sysClr"))
    if sys_clr is not None:
        last = sys_clr.get("lastClr")
        if last:
            return f"#{last}"
    return None


# ============================================================
# Fill / Color Parsing
# ============================================================
def parse_fill(fill_container: ET.Element | None) -> dict:
    if fill_container is None:
        return {"type": "none"}
    if fill_container.find(q("a:noFill")) is not None:
        return {"type": "none"}
    solid = fill_container.find(q("a:solidFill"))
    if solid is not None:
        return {"type": "solid", "color": resolve_color(solid)}
    grad = fill_container.find(q("a:gradFill"))
    if grad is not None:
        stops = []
        gs_lst = grad.find(q("a:gsLst"))
        if gs_lst is not None:
            for gs in gs_lst.findall(q("a:gs")):
                stops.append({"position": gs.get("pos"), "color": resolve_color(gs)})
        return {"type": "gradient", "stops": stops}
    if fill_container.tag.split("}")[-1] == "blipFill":
        blip_elem = fill_container.find(q("a:blip"))
        embed = blip_elem.get(q("r:embed")) if blip_elem is not None else None
        return {"type": "image", "embedRId": embed}
    blip = fill_container.find(q("a:blipFill"))
    if blip is not None:
        blip_elem = blip.find(q("a:blip"))
        embed = blip_elem.get(q("r:embed")) if blip_elem is not None else None
        return {"type": "image", "embedRId": embed}
    return {"type": "unknown"}


def parse_run_color(r_pr: ET.Element | None) -> str | None:
    if r_pr is None:
        return None
    solid = r_pr.find(q("a:solidFill"))
    if solid is not None:
        return resolve_color(solid)
    return None


def parse_text_body(text_body: ET.Element | None) -> dict | None:
    if text_body is None:
        return None
    full_texts = []
    first_color: str | None = None
    for para in text_body.findall(q("a:p")):
        para_text = ""
        for run in para.findall(q("a:r")):
            t = run.find(q("a:t"))
            if t is not None:
                content = "".join(t.itertext())
                para_text += content
                if first_color is None and content.strip():
                    first_color = parse_run_color(run.find(q("a:rPr")))
        if first_color is None:
            p_pr = para.find(q("a:pPr"))
            if p_pr is not None:
                first_color = parse_run_color(p_pr.find(q("a:defRPr")))
        if para_text.strip():
            full_texts.append(para_text.strip())
    full = "\n".join(full_texts)
    if not full:
        return None
    result: dict = {"fullText": full}
    if first_color:
        result["color"] = first_color
    return result


def parse_transform(xfrm: ET.Element | None) -> dict:
    result = {}
    if xfrm is None:
        return result
    off = xfrm.find(q("a:off"))
    ext = xfrm.find(q("a:ext"))
    if off is not None:
        result["x"] = emu_to_pt(off.get("x"))
        result["y"] = emu_to_pt(off.get("y"))
    if ext is not None:
        result["width"] = emu_to_pt(ext.get("cx"))
        result["height"] = emu_to_pt(ext.get("cy"))
    return result


# ============================================================
# Shape / Slide Parsing
# ============================================================
def parse_shape(shape: ET.Element) -> dict | None:
    tag = shape.tag.split("}")[-1]
    if tag not in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
        return None

    result: dict = {"kind": tag}

    xfrm = shape.find(q("p:spPr/a:xfrm"))
    if xfrm is None:
        xfrm = shape.find(q("p:grpSpPr/a:xfrm"))
    result.update(parse_transform(xfrm))

    blip_fill = shape.find(q("p:blipFill"))
    if blip_fill is not None:
        result["fill"] = parse_fill(blip_fill)
    else:
        sp_pr = shape.find(q("p:spPr"))
        if sp_pr is None:
            sp_pr = shape.find(q("p:grpSpPr"))
        if sp_pr is not None:
            result["fill"] = parse_fill(sp_pr)

    tx_body = shape.find(q("p:txBody"))
    if tx_body is not None:
        text = parse_text_body(tx_body)
        if text is not None:
            result["text"] = text

    return result


def parse_background(bg_elem: ET.Element | None) -> dict | None:
    if bg_elem is None:
        return None
    bg_pr = bg_elem.find(q("p:bgPr"))
    if bg_pr is not None:
        return parse_fill(bg_pr)
    bg_ref = bg_elem.find(q("p:bgRef"))
    if bg_ref is not None:
        return {"type": "reference", "idx": bg_ref.get("idx")}
    return None


def parse_slide(slide_path: Path, rels: dict) -> dict:
    root = load_xml(slide_path)
    c_sld = root.find(q("p:cSld")) if root is not None else None
    if c_sld is None:
        return {"background": None, "shapes": [], "media": []}

    bg = parse_background(c_sld.find(q("p:bg")))
    if bg and bg.get("type") == "image" and bg.get("embedRId"):
        rel = rels.get(bg["embedRId"])
        if rel:
            bg["target"] = rel.get("Target", "")

    shapes = []
    sp_tree = c_sld.find(q("p:spTree"))
    if sp_tree is not None:
        for child in sp_tree:
            shape = parse_shape(child)
            if shape is not None:
                shapes.append(shape)

    for shape in shapes:
        if shape.get("kind") == "pic" and shape.get("fill", {}).get("type") == "image":
            embed_r_id = shape["fill"].get("embedRId")
            if embed_r_id:
                rel = rels.get(embed_r_id)
                if rel:
                    shape["fill"]["target"] = rel.get("Target", "")
                    shape["fill"].pop("embedRId", None)

    media = []
    for rel in rels.values():
        target = rel.get("Target", "")
        rel_type = rel.get("Type", "")
        if "/image" in rel_type:
            media.append({
                "target": target,
                "type": Path(target).suffix.lower().lstrip(".") or "unknown",
            })

    return {"background": bg, "shapes": shapes, "media": media}


def parse_theme(theme_path: Path) -> dict:
    root = load_xml(theme_path)
    if root is None:
        return {"colorScheme": {}, "fonts": {}}

    color_scheme = {}
    cs = root.find(f".//{q('a:clrScheme')}")
    if cs is not None:
        for child in cs:
            name = child.tag.split("}")[-1]
            if name == "extLst":
                continue
            color_scheme[name] = resolve_color(child)

    fonts = {}
    fs = root.find(f".//{q('a:fontScheme')}")
    if fs is not None:
        major = fs.find(q("a:majorFont"))
        if major is not None:
            latin = major.find(q("a:latin"))
            ea = major.find(q("a:ea"))
            fonts["heading"] = latin.get("typeface") if latin is not None else None
            fonts["eastAsian"] = ea.get("typeface") if ea is not None else None
        minor = fs.find(q("a:minorFont"))
        if minor is not None:
            latin = minor.find(q("a:latin"))
            fonts["body"] = latin.get("typeface") if latin is not None else None

    return {"colorScheme": color_scheme, "fonts": fonts}


# ============================================================
# Extract PPTX Structure
# ============================================================
def extract(pptx_path: Path, resource_dir: Path | None = None) -> dict:
    """Extract PPTX structure into a Python dict.

    If resource_dir is provided, also extracts all media files into that directory.

    Returns:
        dict with keys: metadata, theme, slides, media_files
        media_files: dict mapping original relative path -> extracted absolute path
    """
    temp_dir = tempfile.mkdtemp(prefix="pptx_")
    pptx_dir = Path(temp_dir)
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            z.extractall(pptx_dir)
    except zipfile.BadZipFile:
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise ValueError(f"not a valid PPTX/ZIP archive: {pptx_path}")

    ppt_dir = pptx_dir / "ppt"
    if not ppt_dir.exists():
        shutil.rmtree(temp_dir, ignore_errors=True)
        raise ValueError("'ppt' subdirectory not found")

    # Page size
    pres_path = ppt_dir / "presentation.xml"
    pres_root = load_xml(pres_path)
    page_size = {}
    if pres_root is not None:
        sld_sz = pres_root.find(q("p:sldSz"))
        if sld_sz is not None:
            page_size = {
                "width": emu_to_pt(sld_sz.get("cx")),
                "height": emu_to_pt(sld_sz.get("cy")),
            }

    theme = parse_theme(ppt_dir / "theme" / "theme1.xml")

    # Parse slides
    slides_dir = ppt_dir / "slides"
    slides = []
    if slides_dir.exists():
        slide_files = sorted(slides_dir.glob("slide*.xml"), key=lambda p: int(re.search(r"\d+", p.stem).group()))
        for slide_file in slide_files:
            rels_path = slides_dir / "_rels" / f"{slide_file.name}.rels"
            slide_rels = load_rels(rels_path)
            slide_data = parse_slide(slide_file, slide_rels)
            slides.append(slide_data)

    # Extract media files to resource_dir
    media_files: dict[str, str] = {}  # original_name -> extracted_path
    media_src_dir = ppt_dir / "media"
    if resource_dir is not None and media_src_dir.exists():
        resource_dir.mkdir(parents=True, exist_ok=True)
        for media_file in sorted(media_src_dir.iterdir()):
            if media_file.is_file():
                dest = resource_dir / media_file.name
                shutil.copy2(media_file, dest)
                media_files[media_file.name] = str(dest)

    shutil.rmtree(temp_dir, ignore_errors=True)

    return {
        "metadata": {"pageSize": page_size},
        "theme": theme,
        "slides": slides,
        "media_files": media_files,
    }


# ============================================================
# Markdown Rendering
# ============================================================
def format_fill(fill: dict | None) -> str:
    if not fill:
        return "无"
    t = fill.get("type", "unknown")
    if t == "none":
        return "无填充"
    if t == "solid":
        return f"纯色 {fill.get('color', '')}"
    if t == "gradient":
        stops = fill.get("stops", [])
        return "渐变 " + ", ".join(f"{s.get('position')}:{s.get('color', '')}" for s in stops)
    if t == "image":
        target = fill.get("target", "")
        return f"图片 `{target}`" if target else "图片（未解析到路径）"
    if t == "reference":
        return f"引用 idx={fill.get('idx', '')}"
    return str(fill)


def format_text(text: dict | None) -> str:
    if not text:
        return ""
    full = text.get("fullText", "")
    return full.replace("|", "\\|").replace("\n", " ")


def render_color_scheme(color_scheme: dict) -> str:
    skip_roles = {"hlink", "folHlink"}
    lines = ["| 角色 | 颜色 | 说明 |", "|------|------|------|"]
    for role, color in color_scheme.items():
        if role in skip_roles:
            continue
        desc = COLOR_ROLE_DESCRIPTIONS.get(role, "")
        lines.append(f"| {role} | {color or '-'} | {desc} |")
    return "\n".join(lines)


def render_fonts(fonts: dict) -> str:
    lines = []
    if fonts.get("heading"):
        lines.append(f"- **标题字体（Heading）**：{fonts['heading']}")
    if fonts.get("body"):
        lines.append(f"- **正文字体（Body）**：{fonts['body']}")
    if fonts.get("eastAsian"):
        lines.append(f"- **东亚字体（EastAsian）**：{fonts['eastAsian']}")
    if not lines:
        lines.append("- （未指定）")
    return "\n".join(lines)


def render_shapes(shapes: list) -> str:
    if not shapes:
        return "（无形状）"
    lines = [
        "| kind | x | y | width | height | fill | text | textColor |",
        "|------|---|---|-------|--------|------|------|-----------|",
    ]
    for sh in shapes:
        kind = sh.get("kind", "")
        if kind == "pic":
            continue
        x = sh.get("x", "")
        y = sh.get("y", "")
        w = sh.get("width", "")
        h = sh.get("height", "")
        fill = format_fill(sh.get("fill"))
        text_obj = sh.get("text") or {}
        text = format_text(text_obj)
        text_color = text_obj.get("color", "")
        lines.append(f"| {kind} | {x} | {y} | {w} | {h} | {fill} | {text} | {text_color} |")
    return "\n".join(lines)


def render_background(bg: dict | None) -> str:
    if not bg:
        return "无（使用主题默认背景）"
    t = bg.get("type", "unknown")
    if t == "image":
        target = bg.get("target", "")
        return f"背景图片： `{target}`" if target else "背景图片（未解析到路径）"
    if t == "solid":
        return f"纯色背景 {bg.get('color', '')}"
    if t == "gradient":
        return f"渐变背景 {format_fill(bg)}"
    if t == "reference":
        return f"引用主题背景 idx={bg.get('idx', '')}"
    return format_fill(bg)


def render_legend() -> str:
    return """## 附录：字段说明

### 全局信息
- **页面尺寸**：幻灯片画布的宽度和高度，单位 pt（1 pt ≈ 1/72 英寸）
- **幻灯片数量**：PPTX 中的总页数

### 主题配色
- **dk1**（dark1）：深色 1，通常用于主要文字颜色
- **lt1**（light1）：浅色 1，通常用于页面背景颜色
- **dk2**（dark2）：深色 2，常用于次级文字或深色强调
- **lt2**（light2）：浅色 2，常用于辅助浅色背景
- **accent1 ~ accent6**：强调色 1~6，常用于品牌色、图表分类色、装饰元素

### 每页背景
- **背景图片**：该页使用图片作为背景时显示其相对路径
- **纯色背景 #RRGGBB**：该页使用纯色背景时显示其十六进制色值
- **渐变背景**：该页使用渐变填充时的简要描述
- **无（使用主题默认背景）**：该页未显式设置背景

### 形状表格字段
- **kind**：形状类型
  - `sp`：普通形状（文本框、矩形、圆形等自定义几何形状）
  - `pic`：图片
  - `grpSp`：组合形状
  - `cxnSp`：连接线
  - `graphicFrame`：图表/表格框架
- **x / y**：形状左上角相对于页面左上角的水平/垂直位置，单位 pt
- **width / height**：形状的宽度和高度，单位 pt
- **fill**：填充样式
  - `无填充`：透明
  - `纯色 #RRGGBB`：纯色填充
  - `渐变 ...`：渐变填充
  - `图片 ../media/...`：图片填充
- **text**：形状中包含的完整文本内容（仅提取文字，不含样式）
- **textColor**：该形状文本的颜色（取首个非空 run 的颜色）
"""


def render(data: dict) -> str:
    parts = []

    parts.append("# PPTX 结构化解析报告")
    parts.append("")

    page_size = data.get("metadata", {}).get("pageSize", {})
    width = page_size.get("width", "")
    height = page_size.get("height", "")
    slide_count = len(data.get("slides", []))

    parts.append("## 全局信息")
    parts.append("")
    parts.append(f"- **页面尺寸**：{width} pt × {height} pt")
    parts.append(f"- **幻灯片数量**：{slide_count}")
    parts.append("")

    theme = data.get("theme", {})
    parts.append("## 主题配色")
    parts.append("")
    parts.append(render_color_scheme(theme.get("colorScheme", {})))
    parts.append("")
    parts.append("## 字体")
    parts.append("")
    parts.append(render_fonts(theme.get("fonts", {})))
    parts.append("")

    for idx, slide in enumerate(data.get("slides", []), start=1):
        parts.append(f"## 第 {idx} 页")
        parts.append("")
        parts.append("### 背景")
        parts.append("")
        parts.append(render_background(slide.get("background")))
        parts.append("")
        parts.append("### 形状")
        parts.append("")
        parts.append(render_shapes(slide.get("shapes", [])))
        parts.append("")

    parts.append(render_legend())
    parts.append("")

    return "\n".join(parts)


def _replace_image_paths(markdown_text: str, resource_base_url: str) -> str:
    """Replace media file references in Markdown with resource_base_url.

    Handles patterns like:
    - `../media/image1.png`
    - `../media/image2.jpeg`
    - 背景图片： `../media/image1.png`

    Replaces with: {resource_base_url}/resource/{filename}
    """
    def replacer(m):
        filename = m.group(1)
        return f"{resource_base_url}/resource/{filename}"

    # Match ../media/xxx.xxx or media/xxx.xxx inside backticks or plain text
    pattern = r"(?:\.\.\/)?media\/([\w\-]+\.\w+)"
    return re.sub(pattern, replacer, markdown_text)


# ============================================================
# Main Entry Points
# ============================================================
def parse_pptx_to_markdown(
    pptx_path: str,
    output_dir: str,
    resource_base_url: str = "",
) -> tuple[str, dict[str, str]]:
    """Parse PPTX to Markdown structure report.

    Args:
        pptx_path: Path to the .pptx file
        output_dir: Directory to extract image resources into (resource/ subdirectory)
        resource_base_url: Optional base URL for image resources.
            If provided, replaces media paths in Markdown with
            {resource_base_url}/resource/{filename}.
            If empty, keeps original relative paths.

    Returns:
        Tuple of (markdown_text, image_files_dict)
        image_files_dict maps filename -> absolute extracted path
    """
    resource_dir = Path(output_dir) / "resource"
    data = extract(Path(pptx_path), resource_dir=resource_dir)
    md_text = render(data)

    if resource_base_url:
        md_text = _replace_image_paths(md_text, resource_base_url)

    media_files = data.get("media_files", {})
    return md_text, media_files


def parse_pptx(pptx_path: str, output_dir: str) -> dict:
    """Backward-compatible entry point.

    Returns dict with 'markdown' key containing the rendered Markdown text,
    and 'media_files' key with extracted image mapping.
    """
    md_text, media_files = parse_pptx_to_markdown(pptx_path, output_dir)
    return {
        "markdown": md_text,
        "media_files": media_files,
    }


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python parse_pptx.py <pptx_path> [output_dir]")
        sys.exit(1)
    pptx_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.dirname(pptx_file)
    md_text, media_files = parse_pptx_to_markdown(pptx_file, out_dir)
    md_path = os.path.join(out_dir, "structure.md")
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(md_text)
    print(f"Markdown report written to: {md_path}")
    print(f"Extracted {len(media_files)} media files")
#!/usr/bin/env python3
"""
PPTX 风格数据提取库
职责：解析 PPTX，输出纯结构化数据 (dict)
不包含任何 prompt 文本，prompt 模板与数据完全分离

用法:
    from scripts.parse_pptx import parse_pptx
    data = parse_pptx("/path/to/file.pptx", "/tmp/output")
"""

import json
import os
import zipfile
import xml.etree.ElementTree as ET
from collections import Counter
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


# ============================================================
# 辅助函数
# ============================================================
def emu_to_inches(emu):
    return round(emu / 914400, 2) if emu else None


def rgb_to_hex(rgb):
    return f"#{str(rgb)}" if rgb else None


# ============================================================
# PPTX 解析
# ============================================================
def extract_theme_info(pptx_path):
    theme = {"colors": {}, "fonts": {}}
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    try:
        with zipfile.ZipFile(pptx_path, "r") as z:
            theme_files = [f for f in z.namelist() if "theme" in f.lower() and f.endswith(".xml")]
            if not theme_files:
                return theme
            root = ET.fromstring(z.read(theme_files[0]))
            clr_scheme = root.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                for child in clr_scheme:
                    tag = child.tag.split("}")[-1]
                    srgb = child.find("a:srgbClr", ns)
                    sys_clr = child.find("a:sysClr", ns)
                    if srgb is not None:
                        theme["colors"][tag] = f"#{srgb.get('val')}"
                    elif sys_clr is not None:
                        theme["colors"][tag] = f"#{sys_clr.get('lastClr', '')}"
            font_scheme = root.find(".//a:fontScheme", ns)
            if font_scheme is not None:
                for role, key in [("majorFont", "major"), ("minorFont", "minor")]:
                    elem = font_scheme.find(f"a:{role}", ns)
                    if elem is not None:
                        latin = elem.find("a:latin", ns)
                        ea = elem.find("a:ea", ns)
                        theme["fonts"][f"{key}_latin"] = latin.get("typeface") if latin is not None else None
                        theme["fonts"][f"{key}_ea"] = ea.get("typeface") if ea is not None else None
    except Exception as e:
        print(f"  [WARN] 主题提取失败: {e}")
    return theme


def collect_shape_style(shape, text_colors, fonts, font_sizes, fill_colors, shape_types):
    shape_types[str(shape.shape_type)] += 1
    try:
        fill = shape.fill
        if fill.type == 1 and fill.fore_color and fill.fore_color.rgb:
            fill_colors[rgb_to_hex(fill.fore_color.rgb)] += 1
    except Exception:
        pass
    try:
        grad = shape._element.find(".//" + qn("a:gradFill"))
        if grad is not None:
            for gs in grad.findall(".//" + qn("a:gs")):
                srgb = gs.find(qn("a:srgbClr"))
                if srgb is not None:
                    fill_colors[f"#{srgb.get('val')}"] += 1
    except Exception:
        pass
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts[run.font.name] += 1
                if run.font.size:
                    font_sizes[run.font.size.pt] += 1
                try:
                    if run.font.color and run.font.color.rgb:
                        text_colors[rgb_to_hex(run.font.color.rgb)] += 1
                except Exception:
                    pass
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            collect_shape_style(sub, text_colors, fonts, font_sizes, fill_colors, shape_types)


def analyze_layout(shapes):
    positions = []
    sw, sh = 9144000, 5143500
    for s in shapes:
        try:
            positions.append({
                "l": round(s.left / sw * 100, 1) if s.left else 0,
                "t": round(s.top / sh * 100, 1) if s.top else 0,
                "w": round(s.width / sw * 100, 1) if s.width else 0,
                "h": round(s.height / sh * 100, 1) if s.height else 0,
            })
        except Exception:
            pass
    if not positions:
        return "empty"
    for p in positions:
        if p["w"] > 80 and p["h"] > 80:
            return "full_bleed"
    if sum(1 for p in positions if p["l"] < 40) >= 2 and sum(1 for p in positions if p["l"] > 50) >= 2:
        return "split_horizontal"
    if sum(1 for p in positions if p["t"] < 30) >= 2 and sum(1 for p in positions if p["t"] > 50) >= 2:
        return "split_vertical"
    if sum(1 for p in positions if 20 < p["l"] < 60) >= len(positions) * 0.6:
        return "centered"
    return "mixed"


def parse_pptx_statistics(prs):
    text_colors, fill_colors, fonts, font_sizes, shape_types = Counter(), Counter(), Counter(), Counter(), Counter()
    slide_summaries = []
    for idx, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        for shape in shapes:
            collect_shape_style(shape, text_colors, fonts, font_sizes, fill_colors, shape_types)
        slide_summaries.append({
            "index": idx + 1,
            "shape_count": len(shapes),
            "has_image": any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes),
            "has_table": any(hasattr(s, "has_table") and s.has_table for s in shapes),
            "has_chart": any(hasattr(s, "has_chart") and s.has_chart for s in shapes),
            "layout_type": analyze_layout(shapes),
        })
    return {
        "text_colors": [{"color": c, "count": n} for c, n in text_colors.most_common(15)],
        "fill_colors": [{"color": c, "count": n} for c, n in fill_colors.most_common(15)],
        "fonts": [{"name": f, "count": n} for f, n in fonts.most_common(10)],
        "font_sizes": [{"size_pt": s, "count": n} for s, n in font_sizes.most_common(15)],
        "shape_types": [{"type": t, "count": n} for t, n in shape_types.most_common(10)],
        "layout_distribution": dict(Counter(s["layout_type"] for s in slide_summaries)),
        "slides_with_images": sum(1 for s in slide_summaries if s["has_image"]),
        "slides_with_tables": sum(1 for s in slide_summaries if s["has_table"]),
        "slides_with_charts": sum(1 for s in slide_summaries if s["has_chart"]),
        "slide_summaries": slide_summaries,
    }


# ============================================================
# 图片提取与主色分析
# ============================================================
def extract_images(prs, images_dir):
    count = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = "png"
                    ct = getattr(image, "content_type", "")
                    if "jpeg" in ct:
                        ext = "jpg"
                    elif "webp" in ct:
                        ext = "webp"
                    with open(os.path.join(images_dir, f"slide{slide_idx+1}_img{count}.{ext}"), "wb") as f:
                        f.write(image.blob)
                    count += 1
                except Exception:
                    pass
    return count


def extract_background_images(pptx_path, bg_images_dir):
    count = 0
    with zipfile.ZipFile(pptx_path, "r") as z:
        for f in z.namelist():
            if "media" in f and f.endswith((".png", ".jpg", ".jpeg")):
                info = z.getinfo(f)
                if info.file_size > 100000:
                    with open(os.path.join(bg_images_dir, os.path.basename(f)), "wb") as out:
                        out.write(z.read(f))
                    count += 1
    return count


def extract_image_colors(image_path, top_n=6):
    try:
        img = Image.open(image_path).convert("RGB")
        img.thumbnail((150, 150))
        pixels = list(img.getdata())
        quantized = [(r // 32 * 32, g // 32 * 32, b // 32 * 32) for r, g, b in pixels]
        return [{"color": f"#{r:02X}{g:02X}{b:02X}", "pct": round(c / len(pixels) * 100, 1)}
                for (r, g, b), c in Counter(quantized).most_common(top_n)]
    except Exception:
        return []


def get_image_dimensions(image_path):
    try:
        with Image.open(image_path) as img:
            return img.size
    except Exception:
        return (0, 0)


def analyze_images(images_dir, bg_images_dir):
    # 背景图
    bg_list = []
    if os.path.exists(bg_images_dir):
        for f in sorted(os.listdir(bg_images_dir)):
            fpath = os.path.join(bg_images_dir, f)
            fsize = os.path.getsize(fpath)
            dims = get_image_dimensions(fpath)
            colors = extract_image_colors(fpath)
            slide_num = None
            for part in f.replace("-", " ").replace(".", " ").split():
                if part.isdigit():
                    slide_num = int(part)
                    break
            bg_list.append({
                "filename": f, "slide": slide_num,
                "dimensions": f"{dims[0]}x{dims[1]}",
                "size_kb": round(fsize / 1024),
                "dominant_colors": colors,
            })
    # 去重合并相同背景图
    unique_bg = {}
    for bg in bg_list:
        key = bg["dimensions"] + str(bg["size_kb"])
        if key not in unique_bg:
            unique_bg[key] = {**bg, "used_in_slides": []}
        unique_bg[key]["used_in_slides"].append(bg["slide"])
    bg_result = []
    for v in unique_bg.values():
        v.pop("slide", None)
        bg_result.append(v)

    # 小图标
    icon_count = 0
    icon_sizes = []
    if os.path.exists(images_dir):
        for f in os.listdir(images_dir):
            fsize = os.path.getsize(os.path.join(images_dir, f))
            if fsize < 5000:
                icon_count += 1
                icon_sizes.append(fsize)

    return {
        "background_images": bg_result,
        "icon_count": icon_count,
        "icon_avg_size_bytes": round(sum(icon_sizes) / len(icon_sizes)) if icon_sizes else 0,
    }


# ============================================================
# 主函数
# ============================================================
def parse_pptx(pptx_path: str, output_dir: str) -> dict:
    """
    解析 PPTX 文件，返回结构化风格数据。

    Args:
        pptx_path: PPTX 文件的绝对路径
        output_dir: 临时图片输出目录（images/ 和 background_images/ 子目录）

    Returns:
        包含 file_info, theme, style_summary, image_analysis 的结构化 dict
    """
    images_dir = os.path.join(output_dir, "images")
    bg_images_dir = os.path.join(output_dir, "background_images")
    os.makedirs(images_dir, exist_ok=True)
    os.makedirs(bg_images_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    w_in = emu_to_inches(prs.slide_width)
    h_in = emu_to_inches(prs.slide_height)
    stats = parse_pptx_statistics(prs)
    theme = extract_theme_info(pptx_path)
    extract_images(prs, images_dir)
    extract_background_images(pptx_path, bg_images_dir)
    image_analysis = analyze_images(images_dir, bg_images_dir)

    return {
        "file_info": {
            "slide_count": len(prs.slides),
            "width_inches": w_in,
            "height_inches": h_in,
            "aspect_ratio": round(w_in / h_in, 2) if h_in else None,
        },
        "theme": theme,
        "style_summary": stats,
        "image_analysis": image_analysis,
    }


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法: python parse_pptx.py <pptx_path> [output_dir]")
        sys.exit(1)
    pptx_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.dirname(pptx_file)
    result = parse_pptx(pptx_file, out_dir)
    json_path = os.path.join(out_dir, "style_data.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"完成: {json_path}")
